#!/usr/bin/env py
# -*- coding: utf-8 -*-
"""
Reference batch driver for /omega:doc-ingest at scale (kg-enhance v2.1).

Converts EVERY ingestible file under the chosen engagement directories to markdown
and stages it under <engagement>/.brain/_ingest_staging/. Loads Marker models ONCE
(not per file) and keeps a resumable manifest so it can stop/restart.

Per-type conversion (each branch independent — a missing converter disables only its type):
  .md   -> direct read
  .pdf  -> Marker PdfConverter (force_ocr=True, languages=ar,en)
  .docx -> Marker DocxConverter if present, else python-docx (paragraphs + tables)
  .pptx -> python-pptx (one '## Slide N' per slide)
  .png/.jpg/.jpeg -> PIL renders a 1-page PDF -> Marker force_ocr

Usage (run from the engagement root, or pass --root):
  py -3.13 brain_batch_ingest.py
  py -3.13 brain_batch_ingest.py --root "D:/path/to/engagement" --dirs 01_Discovery,05_Deliverables_Final
  py -3.13 brain_batch_ingest.py --smoke          # one file of each type
  py -3.13 brain_batch_ingest.py --types md,docx  # subset
  py -3.13 brain_batch_ingest.py --limit 20

Prereqs: pip install marker-pdf python-docx python-pptx Pillow
Note: Arabic OCR is slow (~20-90 s/page); a large PDF corpus can run for hours — run in background.
"""
import os, sys, json, hashlib, time, argparse
from pathlib import Path

DEFAULT_DIRS = ["01_Discovery", "02_Analysis", "03_Recommendations",
                "04_Implementation", "05_Deliverables_Final", "06_Client_Communications"]
DOC_EXT = {".pdf", ".docx", ".pptx"}
IMG_EXT = {".png", ".jpg", ".jpeg"}
MD_EXT = {".md"}
TARGET_EXT = DOC_EXT | IMG_EXT | MD_EXT

_M = {}
def get_marker():
    if not _M:
        from marker.converters.pdf import PdfConverter
        from marker.models import create_model_dict
        from marker.config.parser import ConfigParser
        from marker.output import text_from_rendered
        _M["models"] = create_model_dict()
        _M["PdfConverter"] = PdfConverter
        _M["ConfigParser"] = ConfigParser
        _M["text_from_rendered"] = text_from_rendered
    return _M

def convert_pdf(path):
    m = get_marker()
    cfg = m["ConfigParser"]({"output_format": "markdown", "languages": "ar,en", "force_ocr": True})
    conv = m["PdfConverter"](config=cfg.generate_config_dict(), artifact_dict=m["models"])
    md, _, _ = m["text_from_rendered"](conv(str(path)))
    return md

def convert_image(path, stage):
    from PIL import Image
    img = Image.open(path).convert("RGB")
    tmp = stage / ("_tmp_" + hashlib.md5(str(path).encode()).hexdigest()[:8] + ".pdf")
    img.save(tmp, "PDF")
    try:
        return convert_pdf(tmp)
    finally:
        try: tmp.unlink()
        except OSError: pass

def convert_docx(path):
    try:
        m = get_marker()
        from marker.converters.docx import DocxConverter
        cfg = m["ConfigParser"]({"output_format": "markdown"})
        conv = DocxConverter(config=cfg.generate_config_dict(), artifact_dict=m["models"])
        md, _, _ = m["text_from_rendered"](conv(str(path)))
        return md
    except Exception:
        import docx
        doc = docx.Document(str(path)); out = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if not t: continue
            s = (para.style.name or "").lower() if para.style else ""
            if s.startswith("heading"):
                d = "".join(c for c in s if c.isdigit())
                out.append("#" * (min(int(d), 6) if d else 2) + " " + t)
            else:
                out.append(t)
        for tbl in doc.tables:
            for row in tbl.rows:
                out.append("| " + " | ".join(c.text.strip().replace("\n", " ") for c in row.cells) + " |")
            out.append("")
        return "\n\n".join(out)

def convert_pptx(path):
    from pptx import Presentation
    prs = Presentation(str(path)); out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t: out.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    out.append("| " + " | ".join(c.text.strip().replace("\n", " ") for c in row.cells) + " |")
    return "\n\n".join(out)

def doc_id_for(title, body):
    sha = hashlib.sha256(body.encode("utf-8")).hexdigest()[:12]
    return "DOC-" + hashlib.md5((title + "|1.0|" + sha).encode()).hexdigest()[:10].upper()

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", default=os.getcwd(), help="engagement root (default: cwd)")
    ap.add_argument("--dirs", default=",".join(DEFAULT_DIRS))
    ap.add_argument("--smoke", action="store_true")
    ap.add_argument("--limit", type=int, default=0)
    ap.add_argument("--types", default="")
    args = ap.parse_args()

    root = Path(args.root).resolve()
    dirs = [d.strip() for d in args.dirs.split(",") if d.strip()]
    types = set(t.strip().lower() for t in args.types.split(",") if t.strip())
    stage = root / ".brain" / "_ingest_staging"
    stage.mkdir(parents=True, exist_ok=True)
    manifest_path = stage / "_manifest.json"
    log_path = stage / "_convert.log"

    def log(msg):
        line = f"{time.strftime('%H:%M:%S')} {msg}"
        print(line, flush=True)
        try:
            with open(log_path, "a", encoding="utf-8") as f:
                f.write(line + "\n")
        except OSError:
            pass

    files = []
    for d in dirs:
        base = root / d
        if not base.exists(): continue
        for p in base.rglob("*"):
            if not p.is_file(): continue
            ext = p.suffix.lower()
            if ext not in TARGET_EXT: continue
            if types and ext.lstrip(".") not in types: continue
            if p.name.startswith("~$"): continue
            files.append(p)
    files.sort()
    log(f"Enumerated {len(files)} files under {root.name} (types={types or 'all'})")

    if args.smoke:
        picked, seen = [], set()
        for p in files:
            e = p.suffix.lower()
            if e not in seen:
                seen.add(e); picked.append(p)
        files = picked

    manifest = {}
    if manifest_path.exists():
        try: manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        except Exception: manifest = {}

    def save():
        manifest_path.write_text(json.dumps(manifest, indent=1, ensure_ascii=False), encoding="utf-8")

    done = skipped = failed = processed = 0
    for p in files:
        rel = str(p.relative_to(root)).replace("\\", "/")
        ext = p.suffix.lower()
        rec = manifest.get(rel)
        if rec and rec.get("status") == "done" and (stage / rec.get("staged", "")).exists():
            skipped += 1; continue
        if args.limit and processed >= args.limit:
            break
        processed += 1
        try:
            if ext in MD_EXT:
                body = p.read_text(encoding="utf-8", errors="replace"); conv = "direct"
            elif ext == ".pdf":
                body = convert_pdf(p); conv = "marker-force_ocr"
            elif ext in IMG_EXT:
                body = convert_image(p, stage); conv = "marker-image-ocr"
            elif ext == ".docx":
                body = convert_docx(p); conv = "docx"
            elif ext == ".pptx":
                body = convert_pptx(p); conv = "python-pptx"
            else:
                continue
            body = body or ""
            did = doc_id_for(p.stem, body)
            staged = f"{did}.md"
            fm = ("---\n"
                  f"doc_id: {did}\ndoc_title: {p.stem}\nsource_path: {rel}\n"
                  f"converter: {conv}\nchars: {len(body)}\ntype: document\n---\n\n")
            (stage / staged).write_text(fm + body, encoding="utf-8")
            manifest[rel] = {"status": "done", "converter": conv, "staged": staged,
                             "doc_id": did, "chars": len(body)}
            done += 1
            log(f"[{done}] {conv:18s} {rel} -> {staged} ({len(body)} chars)")
        except Exception as e:
            failed += 1
            manifest[rel] = {"status": "error", "error": f"{type(e).__name__}: {e}"}
            log(f"ERROR {rel}: {type(e).__name__}: {e}")
        if processed % 5 == 0:
            save()
    save()
    log(f"=== COMPLETE targets={len(files)} converted={done} skipped={skipped} failed={failed} ===")

if __name__ == "__main__":
    main()
